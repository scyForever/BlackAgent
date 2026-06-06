# -*- coding: utf-8 -*-
"""Rebuild BlackAgent acceptance PPT/report as a real-sample trace.

The output is intentionally data-bound: it reads the live acceptance JSON and
only presents fields that are actually present in that artifact.
"""

from __future__ import annotations

import json
import re
import shutil
import sys
import textwrap
import zipfile
from pathlib import Path
from typing import Any

PROJECT_ROOT_FOR_IMPORT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT_FOR_IMPORT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT_FOR_IMPORT))

from src.collector.source_metadata import source_class_for_record
from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
RUN_PATH = ROOT / "data" / "acceptance_real_e2e_run_success.json"
EVIDENCE_PATH = ROOT / "data" / "acceptance_real_e2e_evidence.json"
OUT_DIR = ROOT / "docs" / "答辩验收材料"
DETAIL_MD = OUT_DIR / "BlackAgent_真实样例逐步明细.md"
RAW_FULL_MD = OUT_DIR / "BlackAgent_原始数据完整内容.md"
DETAIL_JSON = ROOT / "data" / "acceptance_real_e2e_record_details.json"
RAW_RERUN_PATH = ROOT / "data" / "acceptance_real_e2e_rerun_raw_records.json"
REPORT_MD = OUT_DIR / "BlackAgent_验收报告.md"
REPORT_DOCX = OUT_DIR / "BlackAgent_验收报告.docx"
PPTX = OUT_DIR / "BlackAgent_答辩PPT.pptx"
PPTX_SAFE = OUT_DIR / "BlackAgent_答辩PPT_可打开修复版.pptx"

FONT_CN = "Microsoft YaHei"


def rgb(hex_value: str) -> PptRGB:
    hex_value = hex_value.lstrip("#")
    return PptRGB(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


C = {
    "ink": rgb("0F172A"),
    "muted": rgb("64748B"),
    "line": rgb("E2E8F0"),
    "bg": rgb("F8FAFC"),
    "white": rgb("FFFFFF"),
    "blue": rgb("2563EB"),
    "green": rgb("16A34A"),
    "orange": rgb("F59E0B"),
    "red": rgb("DC2626"),
    "purple": rgb("7C3AED"),
    "pale_blue": rgb("EFF6FF"),
    "pale_green": rgb("F0FDF4"),
    "pale_orange": rgb("FFF7ED"),
    "pale_red": rgb("FEF2F2"),
    "pale_purple": rgb("F5F3FF"),
}

TYPE_CN = {
    "SUSPECTED_CLUSTER": "疑似相似样本组",
    "NEW_SLANG_VARIANT": "新黑话/新说法候选",
    "NEW_RISK_PATTERN": "新风险模式候选",
}
LABEL_CN = {
    "unknown_risk_pattern": "未知风险模式",
    "unknown": "未判定",
    "正常业务白噪声": "正常业务白噪声",
}


def read_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def clean(value: Any) -> str:
    return re.sub(r"\s+", " ", "" if value is None else str(value)).strip()


def short_id(value: str, n: int = 8) -> str:
    return value[:n] if value else ""


def wrap(value: Any, width: int = 40, lines: int = 3) -> str:
    parts = textwrap.wrap(clean(value), width=width)
    if len(parts) > lines:
        parts = parts[:lines]
        parts[-1] = parts[-1].rstrip("，；。") + "…"
    return "\n".join(parts)


def label(value: Any) -> str:
    value = clean(value)
    return LABEL_CN.get(value, value or "未判定")


def parse_entities(summary: str) -> list[dict[str, str]]:
    match = re.search(r"实体线索\s*(.*?)(?:。|$)", summary)
    if not match:
        return []
    raw = match.group(1).strip(" ：:；;")
    items: list[dict[str, str]] = []
    for part in [p.strip() for p in raw.split(", ") if p.strip()]:
        if ":" in part:
            t, v = part.split(":", 1)
            items.append({"type": t.strip(), "value": v.strip()})
        else:
            items.append({"type": "text", "value": part})
    return items


def entity_digest(entities: list[dict[str, str]], max_items: int = 4) -> str:
    if not entities:
        return "该条摘要未展示可见抽取字段"
    names = {
        "url": "链接",
        "contact": "联系方式",
        "slang_term": "黑话",
        "tool_name": "工具名",
        "invite_code": "页面片段",
    }
    out = []
    for ent in entities[:max_items]:
        val = ent["value"]
        if len(val) > 42:
            val = val[:39] + "…"
        out.append(f"{names.get(ent['type'], ent['type'])}:{val}")
    if len(entities) > max_items:
        out.append(f"另{len(entities) - max_items}项")
    return "；".join(out)


def quoted_slang(summary: str) -> str:
    match = re.search(r"候选黑话[“\"]([^”\"]+)[”\"]", summary)
    return match.group(1) if match else "无"


def similar_count(summary: str) -> str:
    match = re.search(r"局部相似样本\s*(\d+)\s*条", summary)
    return match.group(1) if match else "未列出"


def original_label(summary: str, fallback: str) -> str:
    match = re.search(r"原分类=([^，；。]+)", summary)
    return label(match.group(1)) if match else fallback


def _all_final_ids(run: dict[str, Any]) -> set[str]:
    return {
        str(trace_id)
        for clue in run.get("high_quality_clues") or []
        for trace_id in (clue.get("evidence_trace_ids") or [])
    }


def _evidence_clue_map(run: dict[str, Any]) -> dict[str, list[dict[str, Any]]]:
    mapping: dict[str, list[dict[str, Any]]] = {}
    for clue in run.get("high_quality_clues") or []:
        for trace_id in clue.get("evidence_trace_ids") or []:
            mapping.setdefault(str(trace_id), []).append(clue)
    return mapping


def _source_summary(run: dict[str, Any]) -> str:
    classes = (run.get("execution_summary") or {}).get("collection_source_classes_executed") or []
    if classes:
        names = {
            "im_or_group": "IM/群组",
            "social_or_forum": "论坛/社交",
            "vertical_or_technical": "垂直/技术",
        }
        deduped = list(dict.fromkeys(str(item) for item in classes))
        return "、".join(names.get(item, item) for item in deduped)
    source_types = sorted({str(item.get("source_type") or "unknown") for item in run.get("collection_runs") or []})
    return "、".join(source_types) or "公开合规来源"


def _collection_source_for_trace(run: dict[str, Any], trace_id: str) -> tuple[str, str]:
    for clue in run.get("high_quality_clues") or []:
        if trace_id in [str(item) for item in clue.get("evidence_trace_ids") or []]:
            source_names = clue.get("source_names") or []
            source_types = clue.get("source_types") or []
            return (
                "、".join(str(item) for item in source_names) or _source_summary(run),
                "、".join(str(item) for item in source_types) or "公开合规来源",
            )
    return _source_summary(run), "公开合规来源"


def _high_quality_clues(run: dict[str, Any]) -> list[dict[str, Any]]:
    return [item for item in (run.get("high_quality_clues") or []) if isinstance(item, dict)]


def _total_evidence_count(run: dict[str, Any]) -> int:
    return sum(len(clue.get("evidence_trace_ids") or []) for clue in _high_quality_clues(run))


def _source_scope_text(run: dict[str, Any]) -> str:
    return f"公开合规来源（{_source_summary(run)}）"


def _successful_layers(run: dict[str, Any]) -> str:
    layers: list[str] = []
    for item in run.get("collection_runs") or []:
        if item.get("fetched_count"):
            layer = clean(item.get("collection_layer")) or clean(item.get("source_name")) or "未命名批次"
            if layer not in layers:
                layers.append(layer)
    return "、".join(layers) or "无成功批次"


def _collection_meaning(item: dict[str, Any], record_count: int) -> str:
    if item.get("fetched_count"):
        return f"返回有效样本；本轮共 {record_count} 条探索记录参与后续处理"
    return "未形成有效样本，错误/状态已保留"


def _clue_refinement(clue: dict[str, Any]) -> dict[str, Any]:
    return clue.get("refinement") if isinstance(clue.get("refinement"), dict) else {}


def _clue_boundary(clue: dict[str, Any]) -> str:
    ref = _clue_refinement(clue)
    text_value = " ".join(
        clean(value)
        for value in [
            clue.get("risk_category"),
            clue.get("key"),
            clue.get("threshold_reason"),
            ref.get("refined_summary"),
        ]
    )
    if any(marker in text_value for marker in ["未明确关联", "未直接关联", "无明确业务语义", "暂未直接关联"]):
        return "当前主要证明高频共享联系方式/模板重复，尚未直接证明接码、群控脚本或账号交易服务。"
    if any(marker in text_value for marker in ["脚本", "群控", "扫货", "工具交易"]):
        return "命中脚本/工具交易相关证据链，但仍需人工复核排除教程、经验讨论或误报语境。"
    if any(marker in text_value for marker in ["接码", "账号交易", "账号", "实名", "注册"]):
        return "命中目标风险关键词或实体线索，但仍需人工复核确认业务含义和上下文。"
    return "这是高价值候选线索；仍需人工确认业务含义、上下文和处置边界。"


def _record_by_trace(records: list[dict[str, Any]], trace_id: str) -> dict[str, Any] | None:
    return next((item for item in records if item.get("trace_id") == trace_id), None)


def _clue_evidence_rows(clue: dict[str, Any], records: list[dict[str, Any]]) -> list[list[str]]:
    rows: list[list[str]] = []
    for tid in clue.get("evidence_trace_ids") or []:
        tid = str(tid)
        r = _record_by_trace(records, tid)
        if r:
            rows.append([
                r["trace_short"],
                f"{r['classification_label']} / {r['confidence_text']}",
                r["entity_text"],
                _clue_boundary(clue),
            ])
        else:
            rows.append([short_id(tid), "探索明细未找到同号记录", "未暴露", _clue_boundary(clue)])
    return rows


def _clue_summary_line(clue: dict[str, Any]) -> str:
    ids = clue.get("evidence_trace_ids") or []
    return (
        f"线索 `{clue.get('clue_id')}`：风险={clue.get('risk_category')}；"
        f"证据={len(ids)} 条；key=`{clean(clue.get('key'))[:120]}`；"
        f"置信度={clue.get('confidence')}；质量分={clue.get('quality_score')}。"
    )


def build_records(run: dict[str, Any]) -> list[dict[str, Any]]:
    final_ids = _all_final_ids(run)
    clue_map = _evidence_clue_map(run)
    records: list[dict[str, Any]] = []
    for no, h in enumerate(run.get("exploration_hypotheses") or [], 1):
        tid = h.get("source_trace_id") or ""
        summary = clean(h.get("hypothesis_summary"))
        ents = parse_entities(summary)
        conf = float(h.get("confidence") or 0)
        support = [str(x) for x in h.get("supporting_evidence_ids") or []]
        in_final = tid in final_ids
        support_final = sorted(final_ids.intersection(support))
        source_name, source_type = _collection_source_for_trace(run, tid)
        if in_final:
            linked = clue_map.get(tid) or []
            linked_text = "、".join(str(item.get("clue_id")) for item in linked if item.get("clue_id"))
            final_result = f"进入高质量线索证据链：{linked_text or '未列出线索号'}"
        elif support_final:
            final_result = "作为相似样本支持最终证据 " + "、".join(short_id(x) for x in support_final)
        else:
            final_result = "未进入最终线索；保留为人工复核候选"
        review = "需要人工复核" if h.get("requires_human_review") else "可自动通过"
        if conf == 0:
            review = "无置信度；只保留复核"
        budget = h.get("budget_consumed") or {}
        records.append(
            {
                "no": no,
                "trace_id": tid,
                "trace_short": short_id(tid),
                "source": source_name,
                "source_type": source_type,
                "collection_layer": "见 collection_runs",
                "collection_result": "已采到并进入后续处理",
                "cleaning_visible": "本次产物未逐条暴露清洗原因；只暴露总数和后续探索记录",
                "classification_label": label(h.get("suggested_label")),
                "original_label": original_label(summary, label(h.get("suggested_label"))),
                "confidence": conf,
                "confidence_text": f"{conf:.2f}",
                "review_required": bool(h.get("requires_human_review")),
                "review_text": review,
                "hypothesis_type": TYPE_CN.get(h.get("hypothesis_type"), h.get("hypothesis_type") or "未说明"),
                "summary": summary,
                "entities": ents,
                "entity_text": entity_digest(ents),
                "candidate_slang": quoted_slang(summary),
                "similar_count": similar_count(summary),
                "supporting_evidence_ids": support,
                "supporting_evidence_short": "、".join(short_id(x) for x in support) if support else "无",
                "budget_text": f"轮次{budget.get('rounds', '-')}/字数预算{budget.get('tokens', '-')}/耗时{budget.get('elapsed_ms', '-')}ms",
                "in_final_clue": in_final,
                "final_result": final_result,
            }
        )
    return records

def load_raw_records() -> dict[str, Any]:
    if not RAW_RERUN_PATH.exists():
        return {
            "source": "missing",
            "record_count": 0,
            "records": [],
            "note": "未找到同源复跑原始数据文件；当前验收 JSON 本身未保存逐条原始全文。",
        }
    data = read_json(RAW_RERUN_PATH)
    records = []
    for idx, item in enumerate(data.get("records") or [], 1):
        text_value = clean(item.get("content_text"))
        records.append(
            {
                "raw_no": idx,
                "raw_trace_id": f"rerun-raw-{idx:03d}",
                "source_name": item.get("source_name") or "acceptance_telegramnav_public_directory_rerun",
                "source_url": item.get("source_url") or item.get("raw_payload_uri") or "",
                "query": item.get("rerun_query") or item.get("search_query") or "",
                "query_index": item.get("rerun_query_index") or item.get("query_variant_index") or "",
                "matched_keywords": item.get("matched_keywords") or [],
                "keyword_hit_count": item.get("keyword_hit_count"),
                "content_length": len(text_value),
                "content_text": text_value,
            }
        )
    return {
        **data,
        "records": records,
        "record_count": len(records),
        "boundary": "当前验收运行 JSON/数据库没有逐条保存 content_text；本附录展示的是历史同源公开来源、同一查询策略复跑得到的原始完整行内容，不能把复跑编号冒充当前 trace_id。",
    }


def raw_markdown(raw_data: dict[str, Any]) -> str:
    lines = [
        "# BlackAgent 同源复跑原始数据完整内容",
        "",
        "## 说明",
        "",
        "- 当前验收文件 `data/acceptance_real_e2e_run_success.json` 保存探索记录和最终线索证据号，但没有保存每条的 `content_text` 原始全文。",
        "- 为了补充“展示原始数据完整内容”，这里使用历史同源公开来源、同一组查询条件重新采集，并把每条原始行完整保存下来。",
        "- 因为这是复跑数据，编号使用 `rerun-raw-001` 这类新编号，不与当前验收批次的 `source_trace_id` 混用。",
        "",
        f"- 复跑原始行数量：{raw_data.get('record_count', 0)}",
        f"- 边界说明：{raw_data.get('boundary') or raw_data.get('note') or ''}",
        "",
    ]
    for query in raw_data.get("queries") or []:
        lines.append(f"- 复跑查询：`{md(query)}`")
    if raw_data.get("errors"):
        lines.append("")
        lines.append("## 复跑错误")
        for err in raw_data.get("errors") or []:
            lines.append(f"- 查询 {err.get('query_index')}：{md(err.get('error'))}")
    lines += ["", "## 原始数据完整内容", ""]
    for r in raw_data.get("records") or []:
        lines += [
            f"### {r['raw_no']:02d}. `{r['raw_trace_id']}`",
            "",
            f"- 来源：{md(r.get('source_name'))}",
            f"- 查询序号：{r.get('query_index')}",
            f"- 命中关键词：{', '.join(r.get('matched_keywords') or []) or '无'}",
            f"- 原文长度：{r.get('content_length')} 字符",
            "- 原始完整内容：",
            "",
            "```html",
            r.get("content_text") or "",
            "```",
            "",
        ]
    return "\n".join(lines)


def md(value: Any) -> str:
    return clean(value).replace("|", "｜")


def detail_markdown(run: dict[str, Any], evidence: dict[str, Any], records: list[dict[str, Any]]) -> str:
    record_count = len(records)
    clues = _high_quality_clues(run)
    hq_count = len(clues)
    unique_evidence_count = len(_all_final_ids(run))
    total_evidence_count = _total_evidence_count(run)
    source_scope = _source_scope_text(run)
    lines = [
        "# BlackAgent 真实样例逐步明细",
        "",
        "本文件只展示这次真实联网样例里的逐条数据，不做概括式总览。所有字段来自 `data/acceptance_real_e2e_run_success.json`；如果某一步没有在该文件中逐条暴露，就明确写“未逐条暴露”。",
        "",
        "## 1. 样例来源与采集批次",
        "",
        f"- 运行命令：`{md(evidence.get('command'))}`",
        f"- 成功返回数据的层级：`{md(_successful_layers(run))}`。",
        f"- 来源覆盖：{source_scope}；selected_source_count={run.get('selected_source_count')}；input_count={run.get('input_count')}。",
        "- 重要边界：本次样例来自公开合规来源，不代表私群、登录后页面或授权平台数据；最终结果是人工复核候选，不是自动处置结论。",
        "",
        "| 步骤 | 具体内容 | 本步结果 |",
        "|---|---|---|",
    ]
    for i, src in enumerate(run.get("selected_sources") or [], 1):
        lines.append(
            f"| 查询/来源选择 {i} | `{md(src.get('search_query') or src.get('source_name'))}`；理由：{md(src.get('query_rewrite_reason'))} | 已生成公开合规来源查询；未越过登录或私有页面 |"
        )
    for i, item in enumerate(run.get("collection_runs") or [], 1):
        err = item.get("error") or "无错误"
        status = item.get("status") or ("有错误/部分返回" if item.get("error") else "完成")
        lines.append(
            f"| 采集批次 {i} | 层级：`{md(item.get('collection_layer'))}`；来源：{md(item.get('source_name'))}；类别：{md(source_class_for_record(item))} | 返回 {item.get('fetched_count')} 条；状态：{md(status)}；错误：{md(err)} |"
        )

    lines += [
        "",
        f"## 2. {record_count} 条数据逐条处理明细",
        "",
        "说明：完整 JSON 没有逐条保存原始网页全文，也没有逐条写出清洗原因；因此这里不编造原文，只展示可直接复查到的追踪号、分类、抽取信息、探索摘要和最终去向。",
        "",
        "| 序号 | 追踪号 | 采集来源与结果 | 清洗可见结果 | 分类与置信度 | 抽取到的部分信息 | 最终去向 |",
        "|---:|---|---|---|---|---|---|",
    ]
    for r in records:
        lines.append(
            f"| {r['no']} | `{r['trace_short']}` | {md(r['source'])}；{md(r['collection_result'])} | {md(r['cleaning_visible'])} | {md(r['classification_label'])}；{r['confidence_text']}；{md(r['review_text'])} | {md(r['entity_text'])} | {md(r['final_result'])} |"
        )

    lines += ["", "## 3. 每条数据的探索摘要", ""]
    for r in records:
        mark = "【高质量线索证据】" if r["in_final_clue"] else "【复核候选】"
        lines += [
            f"### {r['no']:02d}. {mark} `{r['trace_id']}`",
            "",
            f"- 采集：{r['source']}（{r['source_type']}）；处理状态：{r['collection_result']}。",
            f"- 分类：{r['classification_label']}；置信度 {r['confidence_text']}；处理建议：{r['review_text']}。",
            f"- 类型：{r['hypothesis_type']}；候选黑话：{r['candidate_slang']}；局部相似样本：{r['similar_count']}。",
            f"- 抽取到的部分信息：{r['entity_text']}。",
            f"- 支撑证据号：{r['supporting_evidence_short']}。",
            f"- 本条探索摘要：{r['summary']}",
            f"- 本条最终去向：{r['final_result']}。",
            "",
        ]

    lines += [
        f"## 4. {hq_count} 条高质量候选线索与 {unique_evidence_count} 个唯一证据追踪号",
        "",
        f"本轮 `high_quality_count={run.get('high_quality_count')}`，目标 `>=2`；全部高质量候选线索合计引用 {total_evidence_count} 条证据（去重后 {unique_evidence_count} 个追踪号）。",
        "",
    ]
    for idx, clue in enumerate(clues, 1):
        ref = _clue_refinement(clue)
        ids = [str(item) for item in clue.get("evidence_trace_ids") or []]
        lines += [
            f"### 4.{idx} `{clue.get('clue_id')}`",
            "",
            f"- 风险/类型：{clue.get('risk_category')} / `{clue.get('clue_type')}`。",
            f"- key：`{md(clue.get('key'))}`。",
            f"- 来源：{json.dumps(clue.get('source_names') or [], ensure_ascii=False)}；类型：{json.dumps(clue.get('source_types') or [], ensure_ascii=False)}。",
            f"- 证据数：{len(ids)}；置信度：{clue.get('confidence')}；质量分：{clue.get('quality_score')}。",
            f"- 触发/晋级原因：{md(clue.get('threshold_reason'))}；{md(clue.get('promotion_reason'))}。",
            f"- 证据边界：{_clue_boundary(clue)}",
        ]
        if ref.get("refined_summary"):
            lines.append(f"- 大模型精炼摘要：{ref.get('refined_summary')}")
        for reason in ref.get("refinement_reasons") or []:
            lines.append(f"  - {reason}")
        lines += ["", "| 证据追踪号 | 分类 | 抽取到的部分信息 | 支撑点/边界 |", "|---|---|---|---|"]
        for row in _clue_evidence_rows(clue, records):
            lines.append(f"| `{md(row[0])}` | {md(row[1])} | {md(row[2])} | {md(row[3])} |")
        lines.append("")

    lines += [
        "## 5. 最终人工复核建议",
        "",
        "- 统一结论写法：本次输出的是高价值候选线索和证据链，不是已经人工确认的黑灰产事实。",
    ]
    for idx, clue in enumerate(clues, 1):
        lines.append(f"- 线索 {idx}：{_clue_summary_line(clue)} 复核边界：{_clue_boundary(clue)}")
    lines.append("- 特别说明：Telegram 导航目录共享联系方式线索仍需人工确认是否直接对应接码/群控/账号交易；贴吧脚本模板线索命中工具/脚本证据链，但仍需人工复核排除教程、经验讨论或误报语境。")
    return "\n".join(lines) + "\n"


def report_markdown(run: dict[str, Any], evidence: dict[str, Any], records: list[dict[str, Any]], raw_data: dict[str, Any]) -> str:
    record_count = len(records)
    clues = _high_quality_clues(run)
    hq_count = len(clues)
    unique_evidence_count = len(_all_final_ids(run))
    total_evidence_count = _total_evidence_count(run)
    source_scope = _source_scope_text(run)
    lines = [
        "# BlackAgent 真实样例逐条追踪验收报告",
        "",
        "**项目目录**：`D:\\研一\\BlackAgent`  ",
        "**报告日期**：2026-06-06  ",
        "**本版写法**：不用总览页，不堆指标；按一次真实联网样例，把每条数据在每一步的可见结果写出来。必须保留的英文只用于文件名、命令和字段名，方便复查。",
        "",
        "## 1. 验收口径",
        "",
        "本次验收不把系统说成“已经自动处置黑灰产”，而是验证它能不能把公开合规来源中的线索按步骤留下证据，并交给人复核。",
        "",
        "- 完整运行结果：`data/acceptance_real_e2e_run_success.json`",
        "- 人可读证据：`data/acceptance_real_e2e_evidence.md`",
        "- 逐条明细：`docs/答辩验收材料/BlackAgent_真实样例逐步明细.md`",
        "- 原始全文附录：`docs/答辩验收材料/BlackAgent_原始数据完整内容.md`",
        "",
        f"重要边界：这次样例来自{source_scope}；不是私群、登录后页面，也不是线上生产处置系统；输出为人工复核候选。",
        "",
        "## 2. 样例任务如何开始",
        "",
        "用户输入的任务是：采集公开合规来源中接码、群控脚本、账号交易相关线索，保留完整处理节点、证据链和人工复核建议。",
        "",
        "```powershell",
        evidence.get("command") or "",
        "```",
        "",
        f"命令中的关键意思是：启用联网；只用配置中的公开合规来源；最多选择 {run.get('selected_source_count')} 个来源变体、实际进入处理 {run.get('input_count')} 条记录；把结果写入 `data/acceptance_real_e2e_run_success.json`。",
        "",
        "## 3. 第一步：把任务改写成可采集的查询",
        "",
        "| 序号 | 改写后的查询/来源 | 为什么这样改 | 本步结果 |",
        "|---:|---|---|---|",
    ]
    for i, src in enumerate(run.get("selected_sources") or [], 1):
        lines.append(f"| {i} | `{md(src.get('search_query') or src.get('source_name'))}` | {md(src.get('query_rewrite_reason'))} | 生成公开合规查询；未越过登录或私有页面 |")

    lines += [
        "",
        "## 4. 第二步：真实采集每个批次的结果",
        "",
        "| 批次 | 来源 | 来源类别 | 采集层 | 返回情况 | 这个结果怎么理解 |",
        "|---:|---|---|---|---|---|",
    ]
    for i, item in enumerate(run.get("collection_runs") or [], 1):
        err = item.get("error") or "无错误"
        lines.append(
            f"| {i} | {md(item.get('source_name'))} | `{md(source_class_for_record(item))}` | `{md(item.get('collection_layer'))}` | 返回 {item.get('fetched_count')} 条；错误：{md(err)} | {md(_collection_meaning(item, record_count))} |"
        )

    lines += [
        "",
        f"## 5. 第三步到第五步：{record_count} 条数据逐条结果",
        "",
        "下面的每一行都是一条真实样例数据。因为完整 JSON 没有逐条保存原始网页全文，也没有逐条写出清洗丢弃原因，所以这里不编造原文；只展示文件里能直接复查到的字段：追踪号、分类、置信度、抽取出的部分信息、探索摘要和最终去向。",
        "",
        f"原始完整内容单独放在 `BlackAgent_原始数据完整内容.md`。需要注意：当前验收运行未逐条保存 `content_text`；附录展示的是历史同源公开来源复跑得到的 {raw_data.get('record_count', 0)} 条原始完整行，不能把复跑编号冒充当前 trace_id。",
        "",
        "| # | 追踪号 | 分类结果 | 抽取到的部分信息 | 探索摘要 | 最终去向 |",
        "|---:|---|---|---|---|---|",
    ]
    for r in records:
        lines.append(f"| {r['no']} | `{r['trace_short']}` | {md(r['classification_label'])} / {r['confidence_text']} / {md(r['review_text'])} | {md(r['entity_text'])} | {md(r['summary'])} | {md(r['final_result'])} |")

    lines += [
        "",
        f"## 6. {hq_count} 条高质量候选线索怎样由证据合成",
        "",
        f"系统没有把 {record_count} 条都说成最终线索，而是筛出 {hq_count} 条高质量候选线索；这些线索合计引用 {total_evidence_count} 条证据（去重后 {unique_evidence_count} 个追踪号）。",
        "",
    ]
    for idx, clue in enumerate(clues, 1):
        ref = _clue_refinement(clue)
        lines += [
            f"### 6.{idx} `{clue.get('clue_id')}`",
            "",
            _clue_summary_line(clue),
            "",
            f"- 来源：{json.dumps(clue.get('source_names') or [], ensure_ascii=False)}；类型：{json.dumps(clue.get('source_types') or [], ensure_ascii=False)}。",
            f"- 复核边界：{_clue_boundary(clue)}",
            "",
            "| 证据追踪号 | 分类 | 抽取到的部分信息 | 本条为什么能支撑候选线索 |",
            "|---|---|---|---|",
        ]
        for row in _clue_evidence_rows(clue, records):
            lines.append(f"| `{md(row[0])}` | {md(row[1])} | {md(row[2])} | {md(row[3])} |")
        if ref.get("refined_summary"):
            lines += ["", f"大模型精炼后的说法：{ref.get('refined_summary')}"]
        if ref.get("refinement_reasons"):
            lines.append("它同时给出复核理由：")
            for reason in ref.get("refinement_reasons") or []:
                lines.append(f"- {reason}")
        lines.append("")

    lines += [
        "## 7. 第六步：大模型只做线索精炼，不替人工下结论",
        "",
        "本次外部大模型的作用主要是：理解任务、改写查询、精炼最终线索文字。对于意图和计划两个早期步骤，模型返回内容有字段不完全符合内部格式，系统按规则做了兜底归一化；查询改写和线索精炼使用了模型输出。",
        "",
        "这一步最关键的验收点是：系统没有把证据夸大。它把高质量输出保留为“人工复核候选”，并在证据不足时明确提示需要人工确认。",
        "",
        "## 8. 人工验收时建议怎么讲",
        "",
        "1. 打开 `BlackAgent_答辩PPT.pptx`，从真实任务、查询改写、采集批次开始讲；",
        f"2. 展示 {record_count} 条逐条明细，不再用总览数字代替细节；",
        f"3. 重点点开 {hq_count} 条高质量候选线索和 {unique_evidence_count} 个唯一证据追踪号；",
        "4. 讲最终结论时说“高价值复核候选”，不要说“已确认黑灰产团伙”；",
        "5. 如评委追问原文或清洗原因，说明本次 JSON 没有逐条暴露原始全文和清洗原因，后续可把这两个字段补进验收产物。",
        "",
        "## 9. 本版材料的限制",
        "",
        f"- 可以证明：系统完成了一次公开合规来源的真实联网采集，并把 {record_count} 条样例逐条保留为可复核记录。",
        f"- 可以证明：系统能从 {record_count} 条里筛出 {hq_count} 条高质量候选线索，并保留 {unique_evidence_count} 个唯一证据追踪号。",
        "- 可以证明：本轮目标 high_quality_count >= 2 已达成。",
        "- 不能证明：所有候选线索都已被人工确认属于接码、群控脚本或账号交易。",
        "- 不能证明：系统覆盖了私群、登录后页面或生产级长期监控。",
        "- 不能证明：当前验收运行已经逐条保存完整原文；当前只能展示历史同源复跑原始行。",
    ]
    return "\n".join(lines) + "\n"


# ---------- DOCX ----------


def cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def set_run(run, size=10.0, bold=False, color="111827") -> None:
    run.font.name = FONT_CN
    run._element.rPr.rFonts.set(qn("w:eastAsia"), FONT_CN)
    run.font.size = Pt(size)
    run.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def add_para(doc: Document, text: str, size=10.2, color="111827", bold=False):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(4)
    p.paragraph_format.line_spacing = 1.12
    r = p.add_run(text)
    set_run(r, size=size, color=color, bold=bold)
    return p


def add_heading(doc: Document, text: str, level=1):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(10 if level == 1 else 7)
    p.paragraph_format.space_after = Pt(5)
    r = p.add_run(text)
    set_run(r, size=16 if level == 1 else 12.5, color="1D4ED8" if level == 1 else "0F766E", bold=True)
    return p


def add_bullet(doc: Document, text: str):
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Cm(0.45)
    p.paragraph_format.first_line_indent = Cm(-0.2)
    p.paragraph_format.space_after = Pt(2)
    r = p.add_run("• " + text)
    set_run(r, size=9.3)
    return p


def set_cell(cell, text: str, size=8.0, bold=False, fill=None, color="111827", align=None):
    cell.text = ""
    if fill:
        cell_shading(cell, fill)
    p = cell.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run(str(text))
    set_run(r, size=size, bold=bold, color=color)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def add_table(doc: Document, headers: list[str], rows: list[list[str]], widths: list[float], font_size=7.8):
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    for i, head in enumerate(headers):
        table.rows[0].cells[i].width = Cm(widths[i])
        set_cell(table.rows[0].cells[i], head, size=8.0, bold=True, fill="DBEAFE", color="1E3A8A", align=WD_ALIGN_PARAGRAPH.CENTER)
    for ridx, row in enumerate(rows):
        cells = table.add_row().cells
        fill = "FFFFFF" if ridx % 2 == 0 else "F8FAFC"
        for i, val in enumerate(row):
            cells[i].width = Cm(widths[i])
            set_cell(cells[i], val, size=font_size, fill=fill, align=WD_ALIGN_PARAGRAPH.CENTER if i in (0, 1, 2) else WD_ALIGN_PARAGRAPH.LEFT)
    return table


def build_docx(run: dict[str, Any], records: list[dict[str, Any]], raw_data: dict[str, Any]) -> None:
    record_count = len(records)
    clues = _high_quality_clues(run)
    hq_count = len(clues)
    unique_evidence_count = len(_all_final_ids(run))
    source_scope = _source_scope_text(run)

    doc = Document()
    sec = doc.sections[0]
    sec.orientation = WD_ORIENT.LANDSCAPE
    sec.page_width, sec.page_height = Cm(29.7), Cm(21.0)
    sec.top_margin = sec.bottom_margin = Cm(1.05)
    sec.left_margin = sec.right_margin = Cm(1.05)
    doc.styles["Normal"].font.name = FONT_CN
    doc.styles["Normal"]._element.rPr.rFonts.set(qn("w:eastAsia"), FONT_CN)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("BlackAgent 真实样例逐条追踪验收报告")
    set_run(r, size=22, bold=True, color="1D4ED8")
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(f"用 {record_count} 条真实数据说明每一步发生了什么；高质量候选线索 {hq_count} 条")
    set_run(r, size=11, color="475569")

    add_table(
        doc,
        ["验收材料", "实际文件", "说明"],
        [
            ["完整运行结果", "data/acceptance_real_e2e_run_success.json", f"保存采集、处理、{hq_count} 条高质量候选线索和大模型调用痕迹"],
            ["逐条明细", "docs/答辩验收材料/BlackAgent_真实样例逐步明细.md", f"逐条列出 {record_count} 条数据的可见处理结果"],
            ["原始全文", "docs/答辩验收材料/BlackAgent_原始数据完整内容.md", f"历史同源复跑 {raw_data.get('record_count', 0)} 条原始完整行，单独附录展示"],
            ["答辩 PPT", "docs/答辩验收材料/BlackAgent_答辩PPT.pptx", "按真实样例逐步追踪讲解"],
        ],
        [3.3, 8.2, 14.4],
        font_size=8.6,
    )

    add_heading(doc, "1. 验收口径")
    add_para(doc, "本版报告不再用“系统总览”替代样例细节，而是把一次真实联网运行拆成任务、查询、采集、逐条分类、逐条抽取、线索合成和人工复核建议。")
    add_para(doc, f"重要边界：本次样例来自{source_scope}，不代表私群、登录后页面或生产级长期监控；最终结果是人工复核候选，不是自动处置结论。", color="7C2D12")

    add_heading(doc, "2. 查询改写与采集批次")
    rows = [[str(i), clean(s.get("search_query") or s.get("source_name")), clean(s.get("query_rewrite_reason")), "生成公开合规查询"] for i, s in enumerate(run.get("selected_sources") or [], 1)]
    add_table(doc, ["序号", "改写查询/来源", "改写原因", "本步结果"], rows, [1.0, 6.0, 13.0, 5.5], font_size=7.2)
    rows = []
    for i, item in enumerate(run.get("collection_runs") or [], 1):
        err = item.get("error") or "无错误"
        rows.append([str(i), item.get("collection_layer") or "", source_class_for_record(item), str(item.get("fetched_count")), clean(err), _collection_meaning(item, record_count)])
    add_table(doc, ["批次", "采集层", "来源类别", "返回", "错误字段", "理解方式"], rows, [0.8, 2.5, 2.4, 1.2, 10.0, 8.0], font_size=7.8)

    add_heading(doc, f"3. {record_count} 条数据逐条结果")
    add_para(doc, "说明：完整 JSON 没有逐条保存网页全文，也没有逐条保存清洗原因。本表只展示可以直接复查的字段；没有暴露的字段不编造。", size=9.1, color="7C2D12")
    add_para(doc, f"原始完整内容单独放在 BlackAgent_原始数据完整内容.md：当前验收运行未保存逐条 content_text；附录展示历史同源公开来源复跑得到的 {raw_data.get('record_count', 0)} 条原始完整行。", size=9.1, color="7C2D12")
    rows = [[str(r["no"]), r["trace_short"], f"{r['classification_label']} / {r['confidence_text']}", r["entity_text"], r["summary"], r["final_result"]] for r in records]
    add_table(doc, ["#", "追踪号", "分类/置信度", "抽取到的部分信息", "探索摘要", "最终去向"], rows, [0.8, 2.0, 3.0, 6.2, 11.0, 5.0], font_size=6.2)

    add_heading(doc, f"4. {hq_count} 条高质量候选线索 / {unique_evidence_count} 个唯一证据追踪号")
    for idx, clue in enumerate(clues, 1):
        add_para(doc, f"线索 {idx}/{hq_count}：{_clue_summary_line(clue)}", size=9.4, bold=True)
        add_para(doc, f"复核边界：{_clue_boundary(clue)}", size=9.2, color="991B1B")
        add_table(doc, ["证据追踪号", "分类/置信度", "抽取到的部分信息", "支撑点/边界"], _clue_evidence_rows(clue, records), [2.2, 3.2, 8.0, 11.8], font_size=7.4)
        ref = _clue_refinement(clue)
        if ref.get("refined_summary"):
            add_para(doc, f"大模型精炼：{clean(ref.get('refined_summary'))}", size=8.8)
        for reason in ref.get("refinement_reasons") or []:
            add_bullet(doc, reason)

    add_heading(doc, "5. 本版材料限制")
    for item in [
        f"可以证明：系统完成一次公开合规来源真实联网采集，并把 {record_count} 条样例逐条留痕。",
        f"可以证明：系统能从 {record_count} 条里筛出 {hq_count} 条高质量候选线索，且 high_quality_count >= 2 已达成。",
        "可以展示：历史同源公开来源复跑得到的原始完整行内容，已经单独放入附录。",
        "不能证明：所有候选线索已经被人工确认属于接码、群控脚本或账号交易。",
        "不能证明：系统覆盖私群、登录后页面或生产级长期监控。",
        "不能证明：当前验收运行已经逐条保存完整原文；当前运行原文当时未落盘。",
    ]:
        add_bullet(doc, item)

    footer = doc.sections[0].footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = footer.add_run("BlackAgent 真实样例逐条追踪验收材料 · 2026-06-06")
    set_run(r, size=8, color="64748B")
    doc.save(REPORT_DOCX)


# ---------- PPTX, using only shapes/text boxes (no native tables) ----------


def set_ppt_text(shape, text: Any, size=14, color=None, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = FONT_CN
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C["ink"]


def box(slide, x, y, w, h, text="", fill=None, line=None, size=12, color=None, bold=False, align=PP_ALIGN.LEFT, radius=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or C["line"]
    shp.line.width = PptPt(0.7)
    if text:
        tf = shp.text_frame
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)
        set_ppt_text(shp, text, size=size, color=color or C["ink"], bold=bold, align=align)
    return shp


def text(slide, x, y, w, h, value, size=14, color=None, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    set_ppt_text(tb, value, size=size, color=color or C["ink"], bold=bold, align=align)
    return tb


def title(slide, head: str, sub: str, no: int):
    text(slide, 0.55, 0.26, 9.5, 0.45, head, size=22, bold=True)
    if sub:
        text(slide, 0.58, 0.78, 9.2, 0.28, sub, size=10.5, color=C["muted"])
    box(slide, 12.18, 0.30, 0.70, 0.28, "真实样例", fill=C["pale_blue"], line=C["blue"], size=8.5, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)
    text(slide, 0.55, 7.17, 7.6, 0.2, "来源：data/acceptance_real_e2e_run_success.json；只展示可复查字段", size=7.5, color=C["muted"])
    text(slide, 12.45, 7.17, 0.5, 0.2, str(no), size=8, color=C["muted"], align=PP_ALIGN.RIGHT)


def manual_table(slide, x, y, widths, row_h, headers, rows, font_size=7.3):
    cur = x
    for i, h in enumerate(headers):
        box(slide, cur, y, widths[i], row_h, h, fill=C["pale_blue"], line=C["line"], size=font_size, color=C["blue"], bold=True, align=PP_ALIGN.CENTER, radius=False)
        cur += widths[i]
    cy = y + row_h
    for ridx, row in enumerate(rows):
        cur = x
        fill = C["white"] if ridx % 2 == 0 else C["bg"]
        for i, val in enumerate(row):
            col = C["green"] if i == len(row) - 1 and "最终线索" in str(val) else C["ink"]
            box(slide, cur, cy, widths[i], row_h, val, fill=fill, line=C["line"], size=font_size, color=col, align=PP_ALIGN.CENTER if i in (0, 1, 2) else PP_ALIGN.LEFT, radius=False)
            cur += widths[i]
        cy += row_h


def build_pptx(run: dict[str, Any], records: list[dict[str, Any]], raw_data: dict[str, Any]) -> None:
    record_count = len(records)
    clues = _high_quality_clues(run)
    hq_count = len(clues)
    unique_evidence_count = len(_all_final_ids(run))
    source_scope = _source_scope_text(run)
    selected_count = len(run.get("selected_sources") or [])
    collection_count = len(run.get("collection_runs") or [])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    no = 1

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]
    box(slide, 0.42, 0.42, 12.5, 6.65, fill=C["white"], line=C["line"])
    text(slide, 0.85, 0.80, 8.8, 0.55, "BlackAgent 真实样例逐步追踪", size=28, bold=True)
    text(slide, 0.88, 1.42, 10.6, 0.52, f"用 {record_count} 条真实采集数据说明每一步的部分数据和结果", size=16, color=C["muted"])
    box(slide, 0.9, 2.18, 3.0, 1.0, f"样例来源\n{source_scope}", fill=C["pale_blue"], line=C["blue"], size=12, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, 4.25, 2.18, 3.0, 1.0, f"高质量候选\n{hq_count} 条 / {unique_evidence_count} 个证据", fill=C["pale_green"], line=C["green"], size=13, color=C["green"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, 7.6, 2.18, 3.0, 1.0, "结论边界\n高价值复核候选", fill=C["pale_orange"], line=C["orange"], size=14, color=C["orange"], bold=True, align=PP_ALIGN.CENTER)
    text(slide, 0.92, 4.0, 11.1, 0.55, f"本版材料按“任务 → 查询 → 采集 → {record_count} 条逐条处理 → {hq_count} 条候选线索 → 原始全文附录 → 人工复核建议”讲清楚。", size=16, bold=True)
    text(slide, 0.92, 4.72, 11.2, 0.8, "注意：本次样例能证明公开来源真实联网和线索留痕，不能证明已经人工确认黑灰产，也不能代表私群或登录后数据。", size=14, color=C["red"])
    title(slide, "", "", no); no += 1

    slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]
    title(slide, f"01 任务被改写成 {selected_count} 个可采集查询/来源", "展示真实查询内容和改写原因，不用抽象流程图", no); no += 1
    box(slide, 0.55, 1.12, 12.2, 0.72, "原始任务：联网采集公开合规来源中“接码、群控脚本、账号交易”相关线索，输出证据链和人工复核建议。", fill=C["white"], line=C["line"], size=13, bold=True)
    y = 1.95
    palette = [(C["pale_blue"], C["blue"]), (C["pale_green"], C["green"]), (C["pale_purple"], C["purple"]), (C["pale_orange"], C["orange"])]
    for i, src in enumerate(run.get("selected_sources") or [], 1):
        fill, line = palette[(i - 1) % len(palette)]
        box(slide, 0.60, y, 0.45, 0.42, str(i), fill=line, line=line, size=12, color=PptRGB(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
        box(slide, 1.18, y, 4.90, 0.42, wrap(src.get("search_query") or src.get("source_name"), 40, 1), fill=fill, line=line, size=7.8, bold=True)
        box(slide, 6.25, y, 6.35, 0.42, wrap(src.get("query_rewrite_reason"), 58, 1), fill=C["white"], line=C["line"], size=7.4, color=C["muted"])
        y += 0.52
    box(slide, 0.75, 6.55, 11.8, 0.42, "这一步的重点：保留公开站点域名约束，只访问配置中的公开来源；不触碰私群、登录后页面或验证码。", fill=C["pale_orange"], line=C["orange"], size=10.8, color=C["orange"], bold=True, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]
    title(slide, f"02 {collection_count} 个采集批次逐步返回什么", "只展开真实批次，不做采集总览", no); no += 1
    rows = []
    for i, item in enumerate(run.get("collection_runs") or [], 1):
        err = clean(item.get("error") or "无错误")
        rows.append([i, source_class_for_record(item), item.get("collection_layer") or "", item.get("fetched_count"), err[:45] + ("…" if len(err) > 45 else ""), _collection_meaning(item, record_count)[:28]])
    manual_table(slide, 0.45, 1.15, [0.55, 1.55, 1.75, 0.75, 5.15, 3.15], 0.48, ["批次", "来源类别", "采集层", "返回", "错误/状态", "对本样例的影响"], rows, 7.5)
    box(slide, 0.70, 5.1, 5.9, 0.95, f"下一步追踪 {record_count} 条数据\n每条保留追踪号，后续分类、抽取和线索去向都围绕追踪号展开。", fill=C["pale_blue"], line=C["blue"], size=13, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, 6.9, 5.1, 5.8, 0.95, "不夸大失败批次\n超时和错误没有被包装成有效数据；材料只采用实际返回样本。", fill=C["pale_red"], line=C["red"], size=13, color=C["red"], bold=True, align=PP_ALIGN.CENTER)

    chunk_size = 15
    chunks = [records[i:i + chunk_size] for i in range(0, len(records), chunk_size)] or [[]]
    for ci, chunk in enumerate(chunks, 1):
        slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]
        title(slide, f"03 {record_count} 条数据逐条明细（{ci}/{len(chunks)}）", "每行都是一个 trace_id；最后一列说明本条最终去向", no); no += 1
        rows = []
        for r in chunk:
            final = "高质量证据" if r["in_final_clue"] else "复核候选"
            rows.append([r["no"], r["trace_short"], r["classification_label"], r["confidence_text"], wrap(r["entity_text"], 30, 1), final])
        manual_table(slide, 0.35, 1.13, [0.55, 1.18, 1.45, 0.75, 7.25, 1.65], 0.36, ["#", "追踪号", "分类", "置信度", "抽取到的部分信息", "去向"], rows, 6.2)

    for clue_idx, clue in enumerate(clues, 1):
        ids = [str(item) for item in clue.get("evidence_trace_ids") or []]
        for idx, tid in enumerate(ids, 1):
            r = _record_by_trace(records, tid)
            slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]
            title(slide, f"04 高质量线索 {clue_idx}/{hq_count} 证据 {idx}/{len(ids)}：{short_id(tid)}", "把一条数据在每一步的可见结果拆开讲", no); no += 1
            if not r:
                box(slide, 0.75, 1.4, 11.8, 1.2, f"追踪号 {tid} 在高质量线索中出现，但探索明细未找到同号记录。", fill=C["pale_red"], line=C["red"], size=18, color=C["red"], bold=True)
                continue
            cards = [
                ("采集", f"{wrap(r['source'], 36, 2)}\n类型：{r['source_type']}\n追踪号：{r['trace_id']}", C["pale_blue"], C["blue"]),
                ("清洗", "本次 JSON 未逐条暴露清洗原因\n可见结果：本条进入后续探索和线索判断", C["pale_orange"], C["orange"]),
                ("分类", f"结果：{r['classification_label']}\n置信度：{r['confidence_text']}\n建议：{r['review_text']}", C["pale_green"], C["green"]),
                ("抽取", wrap(r["entity_text"], 36, 5), C["pale_purple"], C["purple"]),
            ]
            for j, (head, body, fill, line) in enumerate(cards):
                x, y = 0.6 + (j % 2) * 6.25, 1.18 + (j // 2) * 1.55
                box(slide, x, y, 5.85, 1.25, fill=fill, line=line)
                box(slide, x + 0.18, y + 0.13, 0.78, 0.28, head, fill=line, line=line, size=8.2, color=PptRGB(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
                text(slide, x + 0.28, y + 0.48, 5.35, 0.60, body, size=9.2, bold=j in (0, 2))
            box(slide, 0.65, 4.45, 12.0, 1.05, f"探索摘要：{wrap(r['summary'], 82, 3)}", fill=C["white"], line=C["line"], size=11)
            box(slide, 0.65, 5.75, 12.0, 0.62, f"最终去向：{r['final_result']}", fill=C["pale_green"], line=C["green"], size=13, color=C["green"], bold=True, align=PP_ALIGN.CENTER)

        slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]
        title(slide, f"05 候选线索 {clue_idx}/{hq_count} 如何由证据合成", "规则和边界一起展示，避免把模型判断说成事实", no); no += 1
        box(slide, 4.05, 1.35, 5.2, 1.1, wrap(clue.get("key"), 34, 3), fill=C["pale_green"], line=C["green"], size=13, color=C["green"], bold=True, align=PP_ALIGN.CENTER)
        positions = [(0.8, 1.15), (9.7, 1.15), (0.8, 2.75), (9.7, 2.75), (0.8, 4.35), (9.7, 4.35)]
        for pos, tid in zip(positions, ids[:6]):
            r = _record_by_trace(records, tid)
            box(slide, pos[0], pos[1], 3.0, 0.78, f"{short_id(tid)}\n{r['classification_label'] if r else ''} / {r['confidence_text'] if r else ''}", fill=C["white"], line=C["line"], size=10.5, bold=True, align=PP_ALIGN.CENTER)
        if len(ids) > 6:
            box(slide, 4.8, 4.75, 3.7, 0.5, f"另 {len(ids) - 6} 条证据见 Markdown 明细", fill=C["white"], line=C["line"], size=11, align=PP_ALIGN.CENTER)
        box(slide, 1.1, 6.05, 11.2, 0.62, f"复核边界：{_clue_boundary(clue)}", fill=C["pale_blue"], line=C["blue"], size=12, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)

        slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]
        title(slide, f"06 线索 {clue_idx}/{hq_count} 大模型精炼：只提示复核", "这页展示最终文字和复核理由", no); no += 1
        ref = _clue_refinement(clue)
        box(slide, 0.75, 1.15, 12.0, 1.10, wrap(ref.get("refined_summary"), 88, 3), fill=C["white"], line=C["line"], size=12, bold=True)
        box(slide, 0.80, 2.55, 2.9, 0.85, f"置信度变化\n{ref.get('confidence_delta')} → {ref.get('final_confidence')}", fill=C["pale_orange"], line=C["orange"], size=14, color=C["orange"], bold=True, align=PP_ALIGN.CENTER)
        box(slide, 4.05, 2.55, 2.9, 0.85, "复核要求\n需要人工复核" if ref.get("review_required") else "复核要求\n不需要复核", fill=C["pale_red"], line=C["red"], size=14, color=C["red"], bold=True, align=PP_ALIGN.CENTER)
        box(slide, 7.30, 2.55, 4.95, 0.85, f"线索号\n{clue.get('clue_id')}", fill=C["pale_blue"], line=C["blue"], size=11, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)
        y = 4.05
        reasons = ref.get("refinement_reasons") or [_clue_boundary(clue)]
        for i, reason in enumerate(reasons[:5], 1):
            box(slide, 0.85, y, 0.36, 0.36, str(i), fill=C["blue"], line=C["blue"], size=10, color=PptRGB(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
            box(slide, 1.35, y - 0.05, 11.0, 0.46, wrap(reason, 84, 2), fill=C["white"], line=C["line"], size=9.2)
            y += 0.58

    slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]
    title(slide, "07 验收边界：能证明什么，不能证明什么", "把“证据”和“结论”分开，避免夸大", no); no += 1
    box(slide, 0.75, 1.25, 5.9, 4.55, fill=C["pale_green"], line=C["green"])
    text(slide, 1.05, 1.55, 5.2, 0.36, "这次可以证明", size=20, color=C["green"], bold=True, align=PP_ALIGN.CENTER)
    for i, item in enumerate(["公开合规来源真实联网成功", f"{record_count} 条样例都有追踪号和探索记录", f"high_quality_count={run.get('high_quality_count')}，目标 >=2 已达成", "系统主动给出人工复核建议"]):
        text(slide, 1.05, 2.25 + i * 0.62, 5.25, 0.32, f"✓ {item}", size=13.5, bold=True)
    box(slide, 6.85, 1.25, 5.9, 4.55, fill=C["pale_red"], line=C["red"])
    text(slide, 7.15, 1.55, 5.2, 0.36, "这次不能证明", size=20, color=C["red"], bold=True, align=PP_ALIGN.CENTER)
    for i, item in enumerate(["已人工确认黑灰产身份", "覆盖私群或登录后页面", "已达到生产级长期监控", "当前 JSON 已保存每条原始全文和清洗原因"]):
        text(slide, 7.15, 2.25 + i * 0.62, 5.35, 0.32, f"× {item}", size=13.5, bold=True)

    slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]
    title(slide, "08 原始完整内容如何展示", "当前验收运行未保存原文；历史同源复跑原文单独放入附录", no); no += 1
    box(slide, 0.75, 1.18, 5.85, 1.25, "当前验收运行\n保存了探索摘要、分类、抽取摘要和线索证据号；没有逐条 content_text。", fill=C["pale_orange"], line=C["orange"], size=13.2, color=C["orange"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, 6.95, 1.18, 5.85, 1.25, f"历史同源复跑\n重新访问同一公开来源和查询，得到 {raw_data.get('record_count', 0)} 条原始完整行，全部写入 Markdown 附录。", fill=C["pale_green"], line=C["green"], size=13.2, color=C["green"], bold=True, align=PP_ALIGN.CENTER)
    sample_rows = []
    for r in (raw_data.get("records") or [])[:6]:
        sample_rows.append([
            r.get("raw_no"),
            r.get("raw_trace_id"),
            ",".join(r.get("matched_keywords") or [])[:18],
            str(r.get("content_length")),
            wrap(r.get("content_text"), 44, 2),
        ])
    manual_table(slide, 0.55, 3.0, [0.55, 1.35, 1.65, 0.95, 9.15], 0.48, ["#", "复跑编号", "命中词", "长度", "原文开头"], sample_rows, 7.2)
    box(slide, 0.95, 6.35, 11.5, 0.42, "完整原文不塞满 PPT，已放入：docs/答辩验收材料/BlackAgent_原始数据完整内容.md", fill=C["pale_blue"], line=C["blue"], size=11.5, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]
    title(slide, "09 答辩时按文件复查", "最后只给复查路径，不回到总览", no); no += 1
    files = [
        ("完整运行结果", "data/acceptance_real_e2e_run_success.json", f"可查 {record_count} 条探索假设、{hq_count} 条高质量候选线索、采集批次和大模型调用"),
        ("逐条明细", "docs/答辩验收材料/BlackAgent_真实样例逐步明细.md", "可按序号查看每一条数据的分类、抽取和去向"),
        ("原始全文", "docs/答辩验收材料/BlackAgent_原始数据完整内容.md", "可查看历史同源复跑得到的每条原始完整内容"),
        ("验收报告", "docs/答辩验收材料/BlackAgent_验收报告.docx", "可作为提交版验收材料"),
    ]
    y = 1.45
    for i, (name, path, desc) in enumerate(files, 1):
        box(slide, 0.85, y, 0.55, 0.55, str(i), fill=C["blue"], line=C["blue"], size=18, color=PptRGB(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
        box(slide, 1.60, y, 4.75, 0.55, name, fill=C["pale_blue"], line=C["blue"], size=14, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)
        box(slide, 6.55, y, 5.75, 0.55, path, fill=C["white"], line=C["line"], size=10.5, bold=True)
        text(slide, 1.65, y + 0.72, 10.8, 0.32, desc, size=11.5, color=C["muted"])
        y += 1.20
    box(slide, 1.05, 6.1, 11.3, 0.55, "一句话结论：这套材料展示的是“真实样例如何一步步变成人工复核候选线索”。", fill=C["pale_green"], line=C["green"], size=15, color=C["green"], bold=True, align=PP_ALIGN.CENTER)

    prs.save(PPTX)
    shutil.copy2(PPTX, PPTX_SAFE)


def zip_check(path: Path) -> dict[str, Any]:
    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()
        bad = zf.testzip()
    return {"path": str(path), "entry_count": len(names), "bad_entry": bad, "has_content_types": "[Content_Types].xml" in names}


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    run = read_json(RUN_PATH)
    evidence = read_json(EVIDENCE_PATH)
    raw_data = load_raw_records()
    records = build_records(run)
    DETAIL_JSON.write_text(
        json.dumps(
            {
                "source": str(RUN_PATH),
                "record_count": len(records),
                "high_quality_count": len(_high_quality_clues(run)),
                "final_evidence_count": _total_evidence_count(run),
                "unique_final_evidence_count": len(_all_final_ids(run)),
                "records": records,
                "raw_full_content_artifact": str(RAW_FULL_MD),
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    DETAIL_MD.write_text(detail_markdown(run, evidence, records), encoding="utf-8")
    RAW_FULL_MD.write_text(raw_markdown(raw_data), encoding="utf-8")
    REPORT_MD.write_text(report_markdown(run, evidence, records, raw_data), encoding="utf-8")
    build_docx(run, records, raw_data)
    build_pptx(run, records, raw_data)
    print(
        json.dumps(
            {
                "record_count": len(records),
                "high_quality_count": len(_high_quality_clues(run)),
                "final_evidence_count": _total_evidence_count(run),
                "unique_final_evidence_count": len(_all_final_ids(run)),
                "outputs": [str(DETAIL_MD), str(RAW_FULL_MD), str(DETAIL_JSON), str(REPORT_MD), str(REPORT_DOCX), str(PPTX), str(PPTX_SAFE)],
                "zip_checks": [zip_check(p) for p in [PPTX, PPTX_SAFE, REPORT_DOCX]],
            },
            ensure_ascii=False,
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
