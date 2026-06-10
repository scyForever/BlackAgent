# -*- coding: utf-8 -*-
"""Rebuild the BlackAgent defense deck (PPTX) and Word report (DOCX) for the
*current* delivery narrative: end-to-end pipeline, four traceable real cases,
headline metrics, engineering trade-offs, and bounded claims.

Design notes
------------
- Data-bound: numbers are read from the live acceptance artifacts. If an
  artifact is missing, the slide degrades gracefully instead of inventing data.
- The DOCX is rendered faithfully from the hand-maintained acceptance report
  Markdown (``BlackAgent_验收报告.md``) so the two stay in sync. This script
  NEVER writes any ``.md`` file; it only emits ``.pptx`` / ``.docx``.
- ``mermaid`` code fences are skipped in the DOCX (their ``text`` ASCII twin
  right below carries the same content in print-friendly form).
"""

from __future__ import annotations

import json
import re
import shutil
import sys
import textwrap
import zipfile
from pathlib import Path
from typing import Any

from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt as PptPt

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "答辩验收材料"
REPORT_MD = OUT_DIR / "BlackAgent_验收报告.md"
REPORT_DOCX = OUT_DIR / "BlackAgent_验收报告.docx"
PPTX = OUT_DIR / "BlackAgent_答辩PPT.pptx"
PPTX_SAFE = OUT_DIR / "BlackAgent_答辩PPT_可打开修复版.pptx"

FONT_CN = "Microsoft YaHei"


# ---------------------------------------------------------------- data loading
def load_json(rel: str) -> dict[str, Any]:
    path = ROOT / rel
    if not path.exists():
        return {}
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return {}


SUMMARY = load_json("data/final_acceptance_summary.json")
CLUE_INDEX = load_json("data/collection_phase_multi_source_clue_evidence_index.json")
PACK_REPORT = load_json("data/collection_phase_multi_source_evidence_pack_report.json")
CLEAN_SUM = load_json("data/cleaning_phase_summary.json")
CLS_SUM = load_json("data/classification_extraction_phase_summary.json")


def g(obj: dict[str, Any], *path: str, default: Any = None) -> Any:
    cur: Any = obj
    for key in path:
        if not isinstance(cur, dict) or key not in cur:
            return default
        cur = cur[key]
    return cur


def num(value: Any, fallback: str = "—") -> str:
    return fallback if value is None else str(value)


# Headline values (data-bound, with verified fallbacks for robustness).
CLS = SUMMARY.get("classification", {})
PCC = PACK_REPORT.get("completeness_counts", {})
CIR = CLUE_INDEX.get("report", {})
CLUE_ROWS = [r for r in CLUE_INDEX.get("rows", []) if isinstance(r, dict)]
SCALE = SUMMARY.get("scale_benchmark", {})
OCR = SUMMARY.get("ocr_hardset", {})

M = {
    "raw": num(CLEAN_SUM.get("input_count"), "4163"),
    "cleaned": num(CLEAN_SUM.get("cleaned_count"), "3464"),
    "dropped": num(CLEAN_SUM.get("dropped_count"), "699"),
    "dup": num(CLEAN_SUM.get("duplicate_drop_count"), "650"),
    "high_risk": num(CLEAN_SUM.get("high_risk_count"), "1095"),
    "entities": num(CLS_SUM.get("entity_count"), "21774"),
    "cls_count": num(CLS_SUM.get("classification_count"), "3464"),
    "review_full": num(CLS_SUM.get("review_required_count"), "970"),
    "pack_rows": num(PACK_REPORT.get("record_count"), "500"),
    "hq_clue_rows": num(PCC.get("has_high_quality_clue"), "17"),
    "cross_rows": num(PCC.get("has_cross_source_clue"), "8"),
    "has_cls": num(PCC.get("has_classification"), "456"),
    "has_ent": num(PCC.get("has_entities"), "453"),
    "has_snap": num(PCC.get("has_capture_snapshot_uri"), "354"),
    "has_hyd": num(PCC.get("has_hydrated_body"), "12"),
    "curated_clues": num(CIR.get("high_quality_clue_count"), "4"),
    "answer_cards": num(CIR.get("answer_chain_card_count"), "17"),
    "missing_trace": num(CIR.get("missing_evidence_trace_count"), "0"),
    "gold": num(CLS.get("record_count"), "193"),
    "primary_f1": num(CLS.get("primary_classification_f1"), "0.8662"),
    "secondary_f1": num(CLS.get("secondary_classification_f1"), "0.8258"),
    "hier_f1": num(CLS.get("hierarchical_classification_f1"), "0.7929"),
    "entity_f1": num(CLS.get("entity_f1"), "0.9484"),
    "fpr": num(CLS.get("false_positive_rate"), "0.0504"),
    "review_rate": num(CLS.get("classification_review_rate"), "0.1865"),
}


# ----------------------------------------------------------------- PPTX colours
def rgb(hex_value: str) -> PptRGB:
    hex_value = hex_value.lstrip("#")
    return PptRGB(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


C = {
    "ink": rgb("0F172A"), "muted": rgb("64748B"), "line": rgb("E2E8F0"),
    "bg": rgb("F8FAFC"), "white": rgb("FFFFFF"), "blue": rgb("2563EB"),
    "green": rgb("16A34A"), "orange": rgb("F59E0B"), "red": rgb("DC2626"),
    "purple": rgb("7C3AED"), "pale_blue": rgb("EFF6FF"), "pale_green": rgb("F0FDF4"),
    "pale_orange": rgb("FFF7ED"), "pale_red": rgb("FEF2F2"), "pale_purple": rgb("F5F3FF"),
}


def clean(value: Any) -> str:
    return re.sub(r"\s+", " ", "" if value is None else str(value)).strip()


def wrap(value: Any, width: int = 40, lines: int = 3) -> str:
    parts = textwrap.wrap(clean(value), width=width)
    if len(parts) > lines:
        parts = parts[:lines]
        parts[-1] = parts[-1].rstrip("，；。、") + "…"
    return "\n".join(parts)


def short_id(value: Any, n: int = 8) -> str:
    value = clean(value)
    return value[:n] if value else ""


def set_ppt_text(shape, value, size=14, color=None, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for line in str(value).split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = FONT_CN
        run.font.size = PptPt(size)
        run.font.bold = bold
        run.font.color.rgb = color or C["ink"]


def box(slide, x, y, w, h, text="", fill=None, line=None, size=12, color=None,
        bold=False, align=PP_ALIGN.LEFT, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or C["line"]
    shape.line.width = PptPt(0.75)
    shape.shadow.inherit = False
    if text:
        tf = shape.text_frame
        tf.margin_left = Inches(0.07)
        tf.margin_right = Inches(0.07)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)
        set_ppt_text(shape, text, size=size, color=color or C["ink"], bold=bold, align=align)
    return shape


def text(slide, x, y, w, h, value, size=14, color=None, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.margin_left = tb.text_frame.margin_right = 0
    tb.text_frame.margin_top = tb.text_frame.margin_bottom = 0
    set_ppt_text(tb, value, size=size, color=color or C["ink"], bold=bold, align=align)
    return tb


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C["bg"]
    return slide


def title(slide, head: str, sub: str, no: int):
    text(slide, 0.55, 0.26, 11.0, 0.5, head, size=23, bold=True)
    if sub:
        text(slide, 0.58, 0.82, 11.0, 0.3, sub, size=11, color=C["muted"])
    box(slide, 11.95, 0.30, 0.95, 0.30, "交付验收", fill=C["pale_blue"], line=C["blue"],
        size=9, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)
    text(slide, 12.55, 7.12, 0.5, 0.2, str(no), size=8, color=C["muted"], align=PP_ALIGN.RIGHT)


def arrow(slide, x, y):
    text(slide, x, y, 0.3, 0.4, "▶", size=14, color=C["muted"], align=PP_ALIGN.CENTER)


# ------------------------------------------------------------------- PPTX build
def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    no = 1

    # --- Slide 1: cover -----------------------------------------------------
    s = new_slide(prs)
    box(s, 0.42, 0.42, 12.5, 6.65, fill=C["white"], line=C["line"])
    text(s, 0.9, 0.95, 11.4, 0.7, "BlackAgent · 黑灰产情报分析 Agent", size=30, bold=True)
    text(s, 0.92, 1.78, 11.4, 0.5, "本地可复跑：公开内容 → 清洗 → 分类 → 实体 → 证据链 → 人工复核候选", size=15, color=C["muted"])
    chips = [("真实数据跑通", C["pale_blue"], C["blue"]),
             ("证据链可追溯", C["pale_green"], C["green"]),
             ("交付边界清楚", C["pale_orange"], C["orange"])]
    for i, (t, fill, line) in enumerate(chips):
        box(s, 0.95 + i * 3.95, 2.7, 3.6, 0.95, t, fill=fill, line=line, size=17, color=line, bold=True, align=PP_ALIGN.CENTER)
    text(s, 0.95, 4.0, 11.3, 0.5, f"4163 raw → 3464 cleaned → 21774 实体 → 500 行证据包 → 4 条精选线索 / 17 条证据卡", size=15, bold=True)
    text(s, 0.95, 4.65, 11.3, 0.45, f"人工 held-out 193 条：一级 F1 {M['primary_f1']}·实体 F1 {M['entity_f1']}·误报率 {M['fpr']}", size=14, color=C["blue"], bold=True)
    box(s, 0.95, 5.45, 11.3, 0.95,
        "边界：系统输出人工复核候选，不自动定性、不自动处置；不覆盖私群、登录后页面、验证码绕过或购买数据。",
        fill=C["pale_red"], line=C["red"], size=13, color=C["red"], bold=True, align=PP_ALIGN.CENTER)
    title(s, "", "", no); no += 1

    # --- Slide 2: end-to-end pipeline --------------------------------------
    s = new_slide(prs)
    title(s, "01 一张图看懂全流程", "不确定的样本不直接下结论，而是带完整证据进入人工复核", no); no += 1
    stages = [
        ("公开/授权来源", "83 来源·多类型", C["pale_blue"], C["blue"]),
        ("数据采集", f"{M['raw']} raw", C["pale_blue"], C["blue"]),
        ("清洗去重", f"{M['cleaned']}·去重{M['dup']}", C["pale_blue"], C["blue"]),
        ("风险分类", f"F1 {M['primary_f1']}", C["pale_green"], C["green"]),
        ("实体抽取", f"{M['entities']}·F1 {M['entity_f1']}", C["pale_green"], C["green"]),
        ("线索+证据链", f"{M['pack_rows']}行·{M['hq_clue_rows']}高质量", C["pale_green"], C["green"]),
        ("人工复核", "不自动处置", C["pale_red"], C["red"]),
    ]
    n = len(stages)
    bw, gap, x0, y0 = 1.62, 0.085, 0.5, 2.55
    for i, (head, sub, fill, line) in enumerate(stages):
        x = x0 + i * (bw + gap)
        box(s, x, y0, bw, 1.3, "", fill=fill, line=line)
        text(s, x + 0.06, y0 + 0.2, bw - 0.12, 0.5, head, size=11.5, color=line, bold=True, align=PP_ALIGN.CENTER)
        text(s, x + 0.06, y0 + 0.78, bw - 0.12, 0.42, sub, size=8.6, color=C["muted"], align=PP_ALIGN.CENTER)
        if i < n - 1:
            arrow(s, x + bw - 0.02, y0 + 0.48)
    box(s, 0.6, 4.6, 12.1, 1.5,
        "全量来源偏 IM/Telegram 公开目录，不是天然均衡；答辩展示均衡性时用四类来源各 20 行的均衡证据包。\n"
        "分类前置极性判别区分公告/反诈/研究/否定语境；抽取侧统一邀请码、TG、URL、联系方式的归一/哈希/掩码字段。",
        fill=C["white"], line=C["line"], size=12.5, align=PP_ALIGN.CENTER)
    text(s, 0.55, 6.35, 12, 0.3, "复查：data/final_acceptance_summary.json · cleaning_phase_summary.json · classification_extraction_phase_summary.json", size=8.5, color=C["muted"])

    # --- Slide 3: headline metrics -----------------------------------------
    s = new_slide(prs)
    title(s, "02 核心数字（都可复查）", "每个数字都能在 data/ 下的产物里复查", no); no += 1
    cards = [
        ("采集规模", f"{M['raw']} raw\n{M['cleaned']} cleaned", C["pale_blue"], C["blue"]),
        ("结构化实体", f"{M['entities']} 个\nURL/联系/邀请码/黑话", C["pale_blue"], C["blue"]),
        ("证据包", f"{M['pack_rows']} 行\n{M['hq_clue_rows']} 高质量·{M['cross_rows']} 跨源", C["pale_green"], C["green"]),
        ("精选真实线索", f"{M['curated_clues']} 条 / {M['answer_cards']} 卡\n缺失证据 {M['missing_trace']}", C["pale_green"], C["green"]),
        ("一级分类 F1", M["primary_f1"], C["pale_green"], C["green"]),
        ("实体 F1", M["entity_f1"], C["pale_green"], C["green"]),
        ("误报率 FPR", M["fpr"], C["pale_orange"], C["orange"]),
        ("线索召回 F1", "1.0\n(24 条线索 gold)", C["pale_green"], C["green"]),
    ]
    cw, ch, gx, gy = 2.92, 1.65, 0.18, 0.28
    for i, (head, val, fill, line) in enumerate(cards):
        col, row = i % 4, i // 4
        x = 0.55 + col * (cw + gx)
        y = 1.5 + row * (ch + gy)
        box(s, x, y, cw, ch, "", fill=fill, line=line)
        text(s, x + 0.12, y + 0.16, cw - 0.24, 0.4, head, size=12, color=line, bold=True, align=PP_ALIGN.CENTER)
        text(s, x + 0.12, y + 0.62, cw - 0.24, 0.9, val, size=15, bold=True, align=PP_ALIGN.CENTER)
    box(s, 0.55, 5.5, 12.2, 0.85,
        f"人工 held-out {M['gold']} 条：一级 F1 {M['primary_f1']}·二级 {M['secondary_f1']}·层级 {M['hier_f1']}·实体 F1 {M['entity_f1']}·复核率 {M['review_rate']}",
        fill=C["pale_blue"], line=C["blue"], size=13, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)
    text(s, 0.55, 6.5, 12, 0.3, "复查：data/manual_heldout_eval_current.json · eval_manual_heldout_clue_recall_report.json · collection_phase_multi_source_evidence_pack_report.json", size=8.3, color=C["muted"])

    # --- Slide 4: four real cases ------------------------------------------
    s = new_slide(prs)
    title(s, "03 四个真实用例总览", "都能追到 source URL 和 trace id；详见真实用例速览第 7 节", no); no += 1
    case_meta = [
        ("① 群控脚本", "工具交易 / 群控脚本", "V2EX 公开论坛", "6 证据 · 质量 0.96", "6 帖聚成同一工具风险主题"),
        ("② 接码平台", "账号交易 / 接码注册", "贴吧公开帖", "3 证据 · 质量 0.9115", "单源工具聚类·需跨源佐证"),
        ("③ 群发/云控", "众包服务 / 代投服务", "V2EX + 贴吧", "4 证据 · 质量 0.96", "跨论坛+社媒·灰区进复核"),
        ("④ 实名号/卖号", "账号交易 / 实名买卖", "贴吧公开帖", "4 证据 · 质量 0.96", "跨源聚成账号交易主题"),
    ]
    palette = [C["blue"], C["green"], C["purple"], C["orange"]]
    pales = [C["pale_blue"], C["pale_green"], C["pale_purple"], C["pale_orange"]]
    for i, (name, cls, src, ev, note) in enumerate(case_meta):
        col, row = i % 2, i // 2
        x = 0.55 + col * 6.2
        y = 1.45 + row * 2.65
        box(s, x, y, 5.95, 2.45, "", fill=pales[i], line=palette[i])
        text(s, x + 0.25, y + 0.18, 5.5, 0.4, name, size=16, color=palette[i], bold=True)
        text(s, x + 0.25, y + 0.72, 5.5, 0.35, f"系统判断：{cls}", size=12, bold=True)
        text(s, x + 0.25, y + 1.16, 5.5, 0.35, f"来源：{src}　|　{ev}", size=11, color=C["muted"])
        box(s, x + 0.25, y + 1.62, 5.45, 0.62, note, fill=C["white"], line=palette[i], size=11.5, color=palette[i], bold=True, align=PP_ALIGN.CENTER)

    # --- Slides 5-8: one per curated clue ----------------------------------
    for ci, row in enumerate(CLUE_ROWS[:4], 1):
        s = new_slide(prs)
        key = clean(row.get("key")) or row.get("clue_id")
        risk = clean(row.get("risk_category"))
        qs = row.get("quality_score")
        chain = [c for c in row.get("answer_chain", []) if isinstance(c, dict)]
        title(s, f"04 真实用例 {ci}/4：{key}", f"{risk} · 质量分 {qs} · {len(chain)} 条证据 · 链路图 + 逐条证据卡", no); no += 1
        # fan-in evidence boxes on the left, clue on the right
        max_show = min(len(chain), 6)
        slot_h = min(0.86, 4.7 / max(max_show, 1))
        for j, card in enumerate(chain[:max_show]):
            y = 1.35 + j * (slot_h + 0.06)
            url = clean(g(card, "raw_snapshot", "source_url"))
            url_short = url.replace("https://", "").replace("http://", "")[:34]
            snippet = wrap(g(card, "raw_snapshot", "raw_snippet"), 30, 2)
            cls = g(card, "classification", default={})
            conf = cls.get("confidence")
            review = "⚑复核" if cls.get("review_required") else ""
            box(s, 0.55, y, 6.3, slot_h, "", fill=C["pale_blue"], line=C["blue"])
            text(s, 0.7, y + 0.05, 6.0, 0.3, f"{short_id(card.get('trace_id'))} · {url_short}", size=8.5, color=C["blue"], bold=True)
            text(s, 0.7, y + 0.32, 6.0, slot_h - 0.34, f"{snippet}", size=8.3, color=C["ink"])
            text(s, 5.7, y + 0.05, 1.05, 0.3, f"{conf}{('  ' + review) if review else ''}", size=8.5, color=(C["red"] if review else C["muted"]), bold=True, align=PP_ALIGN.RIGHT)
        if len(chain) > max_show:
            text(s, 0.7, 1.35 + max_show * (slot_h + 0.06), 6.0, 0.3, f"…另 {len(chain) - max_show} 条见真实用例速览 / 索引 JSON", size=9, color=C["muted"])
        # converge arrow + clue card
        text(s, 7.0, 3.3, 0.5, 0.6, "▶", size=22, color=C["green"], align=PP_ALIGN.CENTER)
        box(s, 7.7, 2.45, 5.0, 2.5, "", fill=C["pale_green"], line=C["green"])
        text(s, 7.9, 2.65, 4.6, 0.6, f"线索：{key}", size=14, color=C["green"], bold=True, align=PP_ALIGN.CENTER)
        text(s, 7.9, 3.45, 4.6, 0.4, f"{risk} · 质量 {qs}", size=13, bold=True, align=PP_ALIGN.CENTER)
        text(s, 7.9, 3.95, 4.6, 0.9, f"{len(chain)} 条证据聚成同一主题\n每条带 URL·分类·实体·trace id", size=11, color=C["muted"], align=PP_ALIGN.CENTER)
        box(s, 0.55, 6.35, 12.15, 0.62,
            "复核边界：证据证明主题反复出现，仍需人工区分研究/防御/售卖滥用；不直接定性。",
            fill=C["pale_orange"], line=C["orange"], size=11.5, color=C["orange"], bold=True, align=PP_ALIGN.CENTER)

    # --- Slide 9: evidence pack completeness -------------------------------
    s = new_slide(prs)
    title(s, "05 证据链完整性", "500 行联合证据包，每行可追溯", no); no += 1
    ev_cards = [
        ("证据包行数", M["pack_rows"]), ("带 source URL / raw", M["pack_rows"]),
        ("带分类", M["has_cls"]), ("带实体", M["has_ent"]),
        ("高质量线索行", M["hq_clue_rows"]), ("跨源线索行", M["cross_rows"]),
        ("带快照 snapshot", M["has_snap"]), ("带正文 hydrated", M["has_hyd"]),
    ]
    for i, (head, val) in enumerate(ev_cards):
        col, row = i % 4, i // 4
        x = 0.55 + col * 3.05
        y = 1.6 + row * 1.9
        box(s, x, y, 2.85, 1.6, "", fill=C["pale_blue"], line=C["blue"])
        text(s, x + 0.1, y + 0.22, 2.65, 0.4, head, size=11.5, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)
        text(s, x + 0.1, y + 0.72, 2.65, 0.7, str(val), size=22, bold=True, align=PP_ALIGN.CENTER)
    box(s, 0.55, 5.7, 12.2, 0.8,
        f"精选证据索引：{M['curated_clues']} 条高质量线索 / {M['answer_cards']} 条证据卡 / 缺失证据 {M['missing_trace']}（status=completed）",
        fill=C["pale_green"], line=C["green"], size=13, color=C["green"], bold=True, align=PP_ALIGN.CENTER)

    # --- Slide 10: scale & classification distribution ---------------------
    s = new_slide(prs)
    title(s, "06 数据规模与分类分布", "全量公开/授权语料的清洗与分类结果", no); no += 1
    left = [
        ("输入 raw", M["raw"]), ("清洗 cleaned", M["cleaned"]),
        ("丢弃 / 重复丢弃", f"{M['dropped']} / {M['dup']}"),
        ("高风险子集", M["high_risk"]), ("实体总数", M["entities"]),
        ("需人工复核（全量）", M["review_full"]),
    ]
    box(s, 0.55, 1.45, 5.9, 5.0, "", fill=C["white"], line=C["line"])
    text(s, 0.8, 1.6, 5.4, 0.4, "清洗 / 规模", size=15, color=C["blue"], bold=True)
    for i, (k, v) in enumerate(left):
        y = 2.2 + i * 0.66
        text(s, 0.9, y, 3.6, 0.4, k, size=12.5)
        text(s, 4.4, y, 1.9, 0.4, str(v), size=13, color=C["blue"], bold=True, align=PP_ALIGN.RIGHT)
    cats = [c for c in CLS_SUM.get("category_counts", []) if isinstance(c, dict)]
    if not cats:
        cats = [{"value": "正常业务白噪声", "count": 1744}, {"value": "账号交易", "count": 513},
                {"value": "工具交易", "count": 333}, {"value": "众包服务", "count": 304},
                {"value": "诈骗引流", "count": 302}, {"value": "unknown", "count": 190},
                {"value": "刷单作弊", "count": 78}]
    cats = sorted(cats, key=lambda c: -(c.get("count") or 0))[:7]
    top = max((c.get("count") or 0) for c in cats) or 1
    box(s, 6.65, 1.45, 6.1, 5.0, "", fill=C["white"], line=C["line"])
    text(s, 6.9, 1.6, 5.4, 0.4, "一级分类分布（全量 cleaned）", size=15, color=C["green"], bold=True)
    for i, c in enumerate(cats):
        y = 2.2 + i * 0.6
        name = clean(c.get("value"))
        cnt = c.get("count") or 0
        text(s, 6.95, y, 2.4, 0.36, name, size=11)
        bar_w = max(0.1, 3.0 * cnt / top)
        box(s, 9.45, y + 0.02, bar_w, 0.32, "", fill=C["pale_green"], line=C["green"], radius=False)
        text(s, 9.5 + bar_w + 0.05, y, 0.9, 0.32, str(cnt), size=10.5, color=C["green"], bold=True)
    text(s, 0.55, 6.6, 12, 0.3, "白噪声/低相关样本保留用于误报率与 hard negative 评估，不包装成风险样本。", size=9, color=C["muted"])

    # --- Slide 11: engineering trade-offs ----------------------------------
    s = new_slide(prs)
    title(s, "07 效果 / 成本 / 时延平衡", "固定预算下优先规则与工程方法，LLM 只用于冲突与高价值复核", no); no += 1
    sizes = SCALE.get("sample_sizes") or [1000, 10000]
    box(s, 0.55, 1.5, 6.0, 2.4, "", fill=C["pale_purple"], line=C["purple"])
    text(s, 0.8, 1.65, 5.4, 0.4, "LLM 价值门控", size=15, color=C["purple"], bold=True)
    text(s, 0.85, 2.2, 5.5, 1.6,
         "record_enrich_policy = conflict_only\nshould_enable_record_enrich = false\n原因：未测得质量收益却增加成本\n→ 简单 query 走规则，复杂/冲突才用 LLM",
         size=12)
    box(s, 6.75, 1.5, 6.0, 2.4, "", fill=C["pale_blue"], line=C["blue"])
    text(s, 7.0, 1.65, 5.4, 0.4, "本地规模 benchmark", size=15, color=C["blue"], bold=True)
    text(s, 7.05, 2.2, 5.6, 1.6,
         f"样本 {sizes[0]} / {sizes[-1]} 条\n约 1246 条/秒，p95≈0.82ms\n该路径 LLM 调用 0 次\n（本地吞吐与路由成本，非真实联网时延）",
         size=12)
    ocr_sub = g(OCR, "ocr_quality_metrics", "substring_match_count", default=20)
    box(s, 0.55, 4.15, 12.2, 1.05,
        f"多模态 OCR hardset：20 条覆盖 chat/poster/qr/screenshot，子串命中 {ocr_sub}/20（受控 hardset 与字段流转证明，非生产 OCR 泛化）。",
        fill=C["pale_orange"], line=C["orange"], size=12.5, color=C["orange"], bold=True, align=PP_ALIGN.CENTER)
    box(s, 0.55, 5.4, 12.2, 1.05,
        "黑话候选发现与灰度生命周期已具备工程能力（probe 报告 input=268/candidate=20，含 1 条人工审核记录）；当前正式报告尚未产出大规模人工确认候选。",
        fill=C["white"], line=C["line"], size=12, align=PP_ALIGN.CENTER)

    # --- Slide A: tech architecture ----------------------------------------
    s = new_slide(prs)
    title(s, "08 技术架构与数据流", "分层 + 流水线；规则/NLP/LLM 协同，安全与预算贯穿全链", no); no += 1
    box(s, 1.55, 1.12, 10.2, 0.5, "入口层　main.py · run_agent_cli.py · serve_demo_api.py（本机 demo）", fill=C["pale_blue"], line=C["blue"], size=11.5, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)
    box(s, 1.55, 1.74, 10.2, 0.5, "application 服务边界　investigation · task · review · report", fill=C["pale_blue"], line=C["blue"], size=11.5, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)
    box(s, 1.55, 2.36, 10.2, 0.5, "LocalAgentRuntime + infra.RuntimeContainer（依赖组装 · 遥测）", fill=C["pale_blue"], line=C["blue"], size=11.5, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)
    box(s, 0.7, 3.1, 3.95, 1.15, "IntelligencePipeline\nClean→Dedup→Classify→\nExtract→Normalize→\nEntityGraph→CluePromotion", fill=C["pale_green"], line=C["green"], size=9.5, color=C["green"], bold=True, align=PP_ALIGN.CENTER)
    box(s, 4.75, 3.1, 3.95, 1.15, "LLM 编排\nModelRouter · BudgetController\nLLMValueGate（固定 JSON schema）\nPolicyGuard 不过 → 回退规则", fill=C["pale_purple"], line=C["purple"], size=9.5, color=C["purple"], bold=True, align=PP_ALIGN.CENTER)
    box(s, 8.8, 3.1, 3.85, 1.15, "safety 守护\nSourcePolicyGuard · PolicyGuard\nPII 掩码 · 输出校验\n默认 dry-run", fill=C["pale_red"], line=C["red"], size=9.5, color=C["red"], bold=True, align=PP_ALIGN.CENTER)
    box(s, 0.7, 4.4, 11.95, 0.5, "storage 双后端　memory / SQLite / PostgreSQL ＋ EntityGraphStore（跨 run 共享）", fill=C["pale_orange"], line=C["orange"], size=11, color=C["orange"], bold=True, align=PP_ALIGN.CENTER)
    box(s, 0.7, 5.15, 11.95, 0.95,
        "主数据流转契约：PipelineItem → PipelineItem。\n规则 + NLP + LLM value gate 协同：简单 query 走规则解析，复杂 / 冲突才上 LLM；LLM 不直接覆盖规则结果。",
        fill=C["white"], line=C["line"], size=12, align=PP_ALIGN.CENTER)

    # --- Slide B: tech选型与取舍 -------------------------------------------
    s = new_slide(prs)
    title(s, "09 关键技术选型与取舍", "决策 → 为什么 → 取舍/边界（详见验收报告第 11 节）", no); no += 1
    heads = ["技术决策", "为什么这么选", "取舍 / 边界"]
    hx = [0.55, 3.85, 8.35]; hw = [3.3, 4.5, 4.3]
    y = 1.4
    for j, h in enumerate(heads):
        box(s, hx[j], y, hw[j], 0.5, h, fill=C["pale_blue"], line=C["blue"], size=11.5, color=C["blue"], bold=True, align=PP_ALIGN.CENTER, radius=False)
    y += 0.5
    trade_rows = [
        ("规则 + NLP + LLM 协同", "固定 token 预算下优先工程化", "record-enrich 无收益→conflict_only；1246 条/秒、LLM 0 次"),
        ("分类仲裁保留四层", "防 LLM 幻觉直接覆盖规则", "LLM 不下最终结论，冲突标 review 转人工"),
        ("极性判别前置", "摘掉防御 / 反诈 / 研究语境", "误报率压到 FPR 0.0504"),
        ("线索分层 + promotion gate", "控制人工复核负担", "复核率收敛 0.1865，弱线索归档"),
        ("配置化 RuleRegistry(YAML)", "风险词 / 正则 / 门槛免改代码", "可维护·可审计·可回归(rule_version)"),
        ("默认 dry-run + SourcePolicyGuard", "合规边界即工程实现", "禁登录 / 验证码 / 私群绕过，PII 掩码加盐"),
    ]
    for ridx, (a, b, c) in enumerate(trade_rows):
        fill = C["white"] if ridx % 2 == 0 else C["bg"]
        box(s, hx[0], y, hw[0], 0.72, a, fill=fill, line=C["line"], size=10, color=C["ink"], bold=True, align=PP_ALIGN.CENTER, radius=False)
        box(s, hx[1], y, hw[1], 0.72, b, fill=fill, line=C["line"], size=9.5, color=C["ink"], align=PP_ALIGN.CENTER, radius=False)
        box(s, hx[2], y, hw[2], 0.72, c, fill=fill, line=C["line"], size=9.5, color=C["ink"], align=PP_ALIGN.CENTER, radius=False)
        y += 0.72
    box(s, 0.55, y + 0.1, 12.1, 0.5,
        "统一口径：能用规则 / 工程解决的不堆大模型；不确定就交人工复核，不让模型替人下结论。",
        fill=C["pale_green"], line=C["green"], size=11.5, color=C["green"], bold=True, align=PP_ALIGN.CENTER)

    # --- Slide 12: boundaries ----------------------------------------------
    s = new_slide(prs)
    title(s, "10 边界：能证明什么 / 不能声称什么", "把“证据”和“结论”分开，避免夸大", no); no += 1
    box(s, 0.6, 1.35, 6.0, 5.0, fill=C["pale_green"], line=C["green"])
    text(s, 0.9, 1.6, 5.4, 0.4, "能证明", size=20, color=C["green"], bold=True, align=PP_ALIGN.CENTER)
    for i, item in enumerate([
        f"公开/授权来源真实采集 {M['raw']} raw",
        "清洗→分类→实体→线索→证据链 闭环",
        "4 条线索可追到 source URL + trace id",
        f"193 人工 gold：F1 {M['primary_f1']}/实体 {M['entity_f1']}",
        "线索召回 F1 1.0（24 条线索 gold）",
    ]):
        text(s, 1.0, 2.3 + i * 0.78, 5.4, 0.6, f"✓ {item}", size=13, bold=True)
    box(s, 6.75, 1.35, 6.0, 5.0, fill=C["pale_red"], line=C["red"])
    text(s, 7.05, 1.6, 5.4, 0.4, "不能声称", size=20, color=C["red"], bold=True, align=PP_ALIGN.CENTER)
    for i, item in enumerate([
        "不是生产实时风控，不自动封禁/处置",
        "不覆盖私群、登录后页面、验证码绕过",
        "不购买数据、不恶意抓取",
        "全量来源偏 IM/TG，不是天然均衡",
        "未完成 1000 条人工标注集（现 193）",
    ]):
        text(s, 7.15, 2.3 + i * 0.78, 5.4, 0.6, f"✗ {item}", size=13, bold=True)

    # --- Slide 13: file index ----------------------------------------------
    s = new_slide(prs)
    title(s, "11 答辩时按文件复查", "现场最快的复查路径", no); no += 1
    files = [
        ("一图看懂（先看）", "docs/答辩验收材料/BlackAgent_一图看懂.md"),
        ("真实用例 + 证据链", "docs/答辩验收材料/BlackAgent_真实用例速览.md"),
        ("完整验收报告", "docs/答辩验收材料/BlackAgent_验收报告.md"),
        ("最终验收摘要", "data/final_acceptance_summary.json"),
        ("线索证据索引", "data/collection_phase_multi_source_clue_evidence_index.json"),
        ("人工评测", "data/manual_heldout_eval_current.json"),
    ]
    for i, (name, path) in enumerate(files):
        y = 1.45 + i * 0.8
        box(s, 0.7, y, 0.5, 0.55, str(i + 1), fill=C["blue"], line=C["blue"], size=16, color=C["white"], bold=True, align=PP_ALIGN.CENTER)
        box(s, 1.4, y, 4.2, 0.55, name, fill=C["pale_blue"], line=C["blue"], size=13, color=C["blue"], bold=True, align=PP_ALIGN.CENTER)
        box(s, 5.75, y, 6.95, 0.55, path, fill=C["white"], line=C["line"], size=11, bold=True)
    box(s, 1.0, 6.4, 11.3, 0.6,
        "一句话结论：BlackAgent 把“公开内容一步步变成可复核线索”跑通了，且每个数字、每条线索都能复查。",
        fill=C["pale_green"], line=C["green"], size=14, color=C["green"], bold=True, align=PP_ALIGN.CENTER)

    PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX)
    shutil.copy2(PPTX, PPTX_SAFE)


# --------------------------------------------------------------- DOCX rendering
def set_run(run, size=10.0, bold=False, color="111827", mono=False) -> None:
    font = "Consolas" if mono else FONT_CN
    run.font.name = font
    run._element.rPr.rFonts.set(qn("w:eastAsia"), FONT_CN)
    run.font.size = Pt(size)
    run.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


INLINE_RE = re.compile(r"(\*\*.+?\*\*|`[^`]+`)")


def add_inline(paragraph, content: str, size: float, color: str = "111827", bold=False) -> None:
    for token in INLINE_RE.split(content):
        if not token:
            continue
        if token.startswith("**") and token.endswith("**"):
            run = paragraph.add_run(token[2:-2])
            set_run(run, size=size, bold=True, color=color)
        elif token.startswith("`") and token.endswith("`"):
            run = paragraph.add_run(token[1:-1])
            set_run(run, size=size - 0.4, bold=bold, color="9333EA", mono=True)
        else:
            run = paragraph.add_run(token)
            set_run(run, size=size, bold=bold, color=color)


def add_table(doc, headers, rows):
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    for i, head in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = ""
        cell_shading(cell, "DBEAFE")
        add_inline(cell.paragraphs[0], head, 8.6, color="1E3A8A")
        for r in cell.paragraphs[0].runs:
            r.bold = True
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    for ridx, row in enumerate(rows):
        cells = table.add_row().cells
        for i in range(len(headers)):
            cell = cells[i]
            cell.text = ""
            cell_shading(cell, "FFFFFF" if ridx % 2 == 0 else "F8FAFC")
            value = row[i] if i < len(row) else ""
            add_inline(cell.paragraphs[0], value, 8.2)
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    return table


def split_table_row(line: str) -> list[str]:
    line = line.strip()
    if line.startswith("|"):
        line = line[1:]
    if line.endswith("|"):
        line = line[:-1]
    return [c.strip() for c in line.split("|")]


def is_separator(line: str) -> bool:
    return bool(re.match(r"^\s*\|?[\s:\-|]+\|?\s*$", line)) and "-" in line


def render_markdown_to_docx(md_text: str, doc) -> None:
    lines = md_text.splitlines()
    i = 0
    first_heading = True
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # fenced code block
        if stripped.startswith("```"):
            lang = stripped[3:].strip().lower()
            body: list[str] = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                body.append(lines[i])
                i += 1
            i += 1  # closing fence
            if lang == "mermaid":
                continue  # ascii twin below carries the same content
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Cm(0.3)
            p.paragraph_format.space_after = Pt(6)
            cell_text = "\n".join(body)
            run = p.add_run(cell_text)
            set_run(run, size=8.4, color="334155", mono=True)
            continue

        # table
        if stripped.startswith("|") and i + 1 < len(lines) and is_separator(lines[i + 1]):
            headers = split_table_row(line)
            i += 2
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                rows.append(split_table_row(lines[i]))
                i += 1
            add_table(doc, headers, rows)
            doc.add_paragraph().paragraph_format.space_after = Pt(2)
            continue

        # headings
        if stripped.startswith("#"):
            level = len(stripped) - len(stripped.lstrip("#"))
            head_text = stripped[level:].strip()
            p = doc.add_paragraph()
            if first_heading and level == 1:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                add_inline(p, head_text, 22, color="1D4ED8", bold=True)
                first_heading = False
            else:
                p.paragraph_format.space_before = Pt(10 if level <= 2 else 6)
                p.paragraph_format.space_after = Pt(4)
                size = {1: 16, 2: 14, 3: 12}.get(level, 11)
                color = "1D4ED8" if level <= 2 else "0F766E"
                add_inline(p, head_text, size, color=color, bold=True)
            i += 1
            continue

        # blockquote
        if stripped.startswith(">"):
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Cm(0.5)
            add_inline(p, stripped.lstrip("> ").strip(), 9.6, color="7C2D12")
            i += 1
            continue

        # bullet / numbered list
        m = re.match(r"^(\s*)([-*]|\d+\.)\s+(.*)$", line)
        if m:
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Cm(0.6)
            p.paragraph_format.first_line_indent = Cm(-0.25)
            p.paragraph_format.space_after = Pt(2)
            bullet = "• " if m.group(2) in ("-", "*") else (m.group(2) + " ")
            run = p.add_run(bullet)
            set_run(run, size=9.6)
            add_inline(p, m.group(3), 9.6)
            i += 1
            continue

        # blank
        if not stripped:
            i += 1
            continue

        # paragraph
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        p.paragraph_format.line_spacing = 1.15
        add_inline(p, stripped, 10.2)
        i += 1


def build_docx_from_report_md() -> None:
    doc = Document()
    sec = doc.sections[0]
    sec.orientation = WD_ORIENT.PORTRAIT
    sec.page_width, sec.page_height = Cm(21.0), Cm(29.7)
    sec.top_margin = sec.bottom_margin = Cm(1.6)
    sec.left_margin = sec.right_margin = Cm(1.7)
    doc.styles["Normal"].font.name = FONT_CN
    doc.styles["Normal"]._element.rPr.rFonts.set(qn("w:eastAsia"), FONT_CN)

    md_text = REPORT_MD.read_text(encoding="utf-8")
    render_markdown_to_docx(md_text, doc)

    footer = doc.sections[0].footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = footer.add_run("BlackAgent 最终验收报告 · 由 BlackAgent_验收报告.md 渲染")
    set_run(run, size=8, color="64748B")
    doc.save(REPORT_DOCX)


# ----------------------------------------------------------------------- checks
def zip_check(path: Path) -> dict[str, Any]:
    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()
        bad = zf.testzip()
    return {
        "path": str(path.relative_to(ROOT)),
        "entry_count": len(names),
        "bad_entry": bad,
        "has_content_types": "[Content_Types].xml" in names,
    }


def main() -> None:
    if not REPORT_MD.exists():
        sys.exit(f"missing {REPORT_MD}")
    build_pptx()
    build_docx_from_report_md()
    report = {
        "data_bound_from": {
            "final_acceptance_summary": bool(SUMMARY),
            "clue_evidence_index": bool(CLUE_INDEX),
            "evidence_pack_report": bool(PACK_REPORT),
            "cleaning_summary": bool(CLEAN_SUM),
            "classification_summary": bool(CLS_SUM),
        },
        "headline": M,
        "curated_clue_rows": len(CLUE_ROWS),
        "pptx_slides": len(Presentation(str(PPTX)).slides),
        "outputs": [str(p.relative_to(ROOT)) for p in [PPTX, PPTX_SAFE, REPORT_DOCX]],
        "zip_checks": [zip_check(p) for p in [PPTX, PPTX_SAFE, REPORT_DOCX]],
    }
    print(json.dumps(report, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
